#!/usr/bin/env python3
"""
PDF to PPTX Converter
=====================
Converts each page of a PDF file into a slide in a PowerPoint presentation.
Each PDF page is rendered as a high-resolution image and placed on a slide.

Usage:
    python pdf_to_pptx.py input.pdf                    # Output: input.pptx
    python pdf_to_pptx.py input.pdf -o output.pptx     # Custom output name
    python pdf_to_pptx.py input.pdf --dpi 300           # Higher quality (default: 200)
"""

import argparse
import io
import sys
import os

try:
    import fitz  # PyMuPDF
except ImportError:
    print("❌ PyMuPDF is not installed. Run: pip install PyMuPDF")
    sys.exit(1)

try:
    from pptx import Presentation
    from pptx.util import Inches, Emu
except ImportError:
    print("❌ python-pptx is not installed. Run: pip install python-pptx")
    sys.exit(1)

try:
    from PIL import Image
except ImportError:
    print("❌ Pillow is not installed. Run: pip install Pillow")
    sys.exit(1)


def pdf_to_pptx(pdf_path: str, output_path: str, dpi: int = 200) -> str:
    """
    Convert a PDF file to a PPTX presentation.

    Each page of the PDF is rendered at the given DPI and placed as a
    full-slide image in the resulting PowerPoint file.

    Args:
        pdf_path:    Path to the input PDF file.
        output_path: Path for the output PPTX file.
        dpi:         Render resolution (default 200 — good balance of quality/size).

    Returns:
        The path to the created PPTX file.
    """
    if not os.path.isfile(pdf_path):
        raise FileNotFoundError(f"PDF file not found: {pdf_path}")

    doc = fitz.open(pdf_path)
    total_pages = len(doc)

    if total_pages == 0:
        raise ValueError("The PDF file has no pages.")

    prs = Presentation()

    # Use a blank slide layout
    blank_layout = prs.slide_layouts[6]  # index 6 is typically the blank layout

    print(f"📄 Converting '{pdf_path}' ({total_pages} page{'s' if total_pages != 1 else ''}) ...")

    for page_num in range(total_pages):
        page = doc[page_num]

        # Render the page at the specified DPI
        zoom = dpi / 72  # PDF default is 72 DPI
        mat = fitz.Matrix(zoom, zoom)
        pix = page.get_pixmap(matrix=mat, alpha=False)

        # Convert pixmap to PIL Image, then to bytes for python-pptx
        img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
        img_stream = io.BytesIO()
        img.save(img_stream, format="PNG")
        img_stream.seek(0)

        # Calculate slide dimensions to match PDF page aspect ratio
        page_width_inches = page.rect.width / 72
        page_height_inches = page.rect.height / 72

        # Set presentation slide size to match the first page
        if page_num == 0:
            prs.slide_width = Inches(page_width_inches)
            prs.slide_height = Inches(page_height_inches)

        slide = prs.slides.add_slide(blank_layout)

        # Scale image to fill the slide while preserving aspect ratio
        slide_w = prs.slide_width
        slide_h = prs.slide_height
        img_aspect = pix.width / pix.height
        slide_aspect = slide_w / slide_h

        if img_aspect > slide_aspect:
            # Image is wider — fit to width
            pic_width = slide_w
            pic_height = int(slide_w / img_aspect)
        else:
            # Image is taller — fit to height
            pic_height = slide_h
            pic_width = int(slide_h * img_aspect)

        left = int((slide_w - pic_width) / 2)
        top = int((slide_h - pic_height) / 2)

        slide.shapes.add_picture(img_stream, left, top, pic_width, pic_height)

        progress = int((page_num + 1) / total_pages * 100)
        print(f"   ✅ Page {page_num + 1}/{total_pages} — {progress}%")

    doc.close()

    prs.save(output_path)
    print(f"\n🎉 Done! Saved to '{output_path}'")
    return output_path


def main():
    parser = argparse.ArgumentParser(
        description="Convert a PDF file to a PowerPoint (PPTX) presentation.",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  python pdf_to_pptx.py report.pdf
  python pdf_to_pptx.py report.pdf -o slides.pptx
  python pdf_to_pptx.py report.pdf --dpi 300
        """,
    )
    parser.add_argument("pdf", help="Path to the input PDF file")
    parser.add_argument(
        "-o", "--output",
        help="Output PPTX file path (default: same name as PDF with .pptx extension)",
    )
    parser.add_argument(
        "--dpi",
        type=int,
        default=200,
        help="Render DPI for PDF pages (default: 200). Higher = better quality but larger file.",
    )

    args = parser.parse_args()

    pdf_path = args.pdf
    if args.output:
        output_path = args.output
    else:
        output_path = os.path.splitext(pdf_path)[0] + ".pptx"

    try:
        pdf_to_pptx(pdf_path, output_path, dpi=args.dpi)
    except FileNotFoundError as e:
        print(f"❌ {e}")
        sys.exit(1)
    except Exception as e:
        print(f"❌ Error: {e}")
        sys.exit(1)


if __name__ == "__main__":
    main()
