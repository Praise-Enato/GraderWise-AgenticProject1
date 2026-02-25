"""
Training Data Loader — Extracts human-graded scores from the BYUMS rubric
spreadsheet and text content from the corresponding PPTX business plans.

Usage:
    python -m backend.src.training_data_loader
"""

import os
import logging
from typing import List, Dict, Optional

logger = logging.getLogger(__name__)

# Path to training data directory
TRAINING_DATA_DIR = os.path.join(os.path.dirname(__file__), "..", "training_data")

from backend.src.business_rubric_templates import BYUMS_RUBRIC

# Row indices from the BYUMS Excel spreadsheet matching the 24 criteria
_excel_rows = [13, 14, 15, 19, 20, 21, 22, 26, 27, 28, 29, 33, 34, 35, 36, 40, 41, 42, 46, 47, 51, 55, 56, 57]

# Dynamically map the spreadsheet row to the exact string criteria from templates
RUBRIC_ROW_MAP = [
    (r, BYUMS_RUBRIC[i].criteria, BYUMS_RUBRIC[i].max_points, BYUMS_RUBRIC[i].criteria.split(" - ")[0])
    for i, r in enumerate(_excel_rows)
]

# Business plan column mapping (0-indexed, columns B through G = indices 1-6)
BUSINESS_COLUMNS = [
    (1, "Jideofor Enterprise", "Logistics"),
    (2, "Princess Enterprises", "Finance"),
    (3, "Mwana Mboka Logistics Sarlu", "Agriculture"),
    (4, "Kalemie Mobile Health Outreach", "Healthcare"),
    (5, "Kachlinks Technologies", "Printing"),
    (6, "Mulla Global Systems", "Solar"),
]

# Map business names to their PPTX filenames
PPTX_FILENAME_MAP = {
    "Jideofor Enterprise": "Jideofor Enterprise.pptx",
    "Princess Enterprises": "Princess Enterprises.pptx",
    "Mwana Mboka Logistics Sarlu": "Mwana Mboka Logistics.pptx",
    "Kalemie Mobile Health Outreach": "Kalemie Mobile Outreach.pptx",
    "Kachlinks Technologies": "Kachlinks Technologies.pptx",
    "Mulla Global Systems": "Mulla Global Syatems.pptx",  # Note: typo in filename
}


def extract_pptx_text(filepath: str) -> str:
    """Extract all text content from a PPTX file."""
    try:
        from pptx import Presentation
        prs = Presentation(filepath)
        all_text = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            slide_texts.append(text)
                # Also extract table content
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = " | ".join(
                            cell.text.strip() for cell in row.cells if cell.text.strip()
                        )
                        if row_text:
                            slide_texts.append(row_text)
            if slide_texts:
                all_text.append(f"[Slide {slide_num}]\n" + "\n".join(slide_texts))
        return "\n\n".join(all_text)
    except Exception as e:
        logger.error(f"Error extracting text from {filepath}: {e}")
        return ""


def load_training_data() -> List[Dict]:
    """
    Load all training data from the BYUMS rubric spreadsheet and PPTX files.

    Returns:
        List of dicts, each representing one graded business plan:
        {
            "business_name": str,
            "sector": str,
            "pptx_text": str (may be empty for image-based plans),
            "grand_total": float,
            "criteria_scores": [
                {
                    "criteria": str,
                    "max_points": float,
                    "awarded_points": float,
                    "category": str,
                }
            ],
            "category_totals": {
                "Video": {"awarded": float, "max": float},
                ...
            }
        }
    """
    import openpyxl

    xlsx_path = os.path.join(TRAINING_DATA_DIR, "BYUMS Business Plan Competition Rubric.xlsx")
    if not os.path.exists(xlsx_path):
        logger.error(f"Rubric spreadsheet not found: {xlsx_path}")
        return []

    wb = openpyxl.load_workbook(xlsx_path, data_only=True)
    ws = wb["Sheet1"]

    # Read all cell values into a 2D array for easier access
    all_rows = []
    for row in ws.iter_rows(min_row=1, max_row=ws.max_row, max_col=8, values_only=True):
        all_rows.append(list(row))

    training_examples = []

    for col_idx, biz_name, sector in BUSINESS_COLUMNS:
        criteria_scores = []
        category_totals = {}

        for row_idx, criteria_name, max_pts, category in RUBRIC_ROW_MAP:
            score = all_rows[row_idx][col_idx]
            try:
                awarded = float(score) if score is not None else 0.0
            except ValueError:
                awarded = 0.0

            criteria_scores.append({
                "criteria": criteria_name,
                "max_points": max_pts,
                "awarded_points": awarded,
                "category": category,
            })

            # Accumulate category totals
            if category not in category_totals:
                category_totals[category] = {"awarded": 0.0, "max": 0.0}
            category_totals[category]["awarded"] += awarded
            category_totals[category]["max"] += max_pts

        grand_total = sum(c["awarded_points"] for c in criteria_scores)

        # Extract PPTX text
        pptx_filename = PPTX_FILENAME_MAP.get(biz_name, "")
        pptx_path = os.path.join(TRAINING_DATA_DIR, pptx_filename)
        pptx_text = ""
        if os.path.exists(pptx_path):
            pptx_text = extract_pptx_text(pptx_path)
        else:
            logger.warning(f"PPTX file not found for {biz_name}: {pptx_path}")

        training_examples.append({
            "business_name": biz_name,
            "sector": sector,
            "pptx_text": pptx_text,
            "grand_total": grand_total,
            "criteria_scores": criteria_scores,
            "category_totals": category_totals,
        })

    logger.info(f"Loaded {len(training_examples)} training examples")
    return training_examples


if __name__ == "__main__":
    logging.basicConfig(level=logging.INFO)
    examples = load_training_data()
    for ex in examples:
        print(f"\n{'='*60}")
        print(f"Business: {ex['business_name']} ({ex['sector']})")
        print(f"Grand Total: {ex['grand_total']}/100")
        print(f"PPTX Text Length: {len(ex['pptx_text'])} chars")
        print(f"Category Totals:")
        for cat, totals in ex["category_totals"].items():
            print(f"  {cat}: {totals['awarded']}/{totals['max']}")
        print(f"Criteria ({len(ex['criteria_scores'])} items):")
        for c in ex["criteria_scores"]:
            print(f"  [{c['category']}] {c['criteria']}: {c['awarded_points']}/{c['max_points']}")
