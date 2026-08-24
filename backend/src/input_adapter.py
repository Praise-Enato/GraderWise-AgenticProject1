"""Input adapters — normalize any submission source to a common shape.

The submission source will change (PDF for the demo; possibly a Google Slides
link later). To avoid painting the pipeline into a corner, everything goes
through an adapter that yields a NormalizedInput{ text, page_images, video_url }.
Swapping PDF -> Slides later is a new adapter, not a pipeline rewrite (same idea
as the model router).

For the demo only the PDF adapter is built. The pure parts (YouTube-link
extraction, adapter routing, the normalized shape) are unit-tested; PDF text and
page-image extraction lazily import pypdf / pypdfium2 and degrade gracefully when
those libraries are absent.
"""
from __future__ import annotations

import re
from dataclasses import dataclass, field
from typing import List, Optional


@dataclass
class NormalizedInput:
    text: str = ""
    page_images: List[bytes] = field(default_factory=list)  # PNG bytes per page (for vision model)
    video_url: Optional[str] = None
    source: str = ""
    notes: List[str] = field(default_factory=list)  # non-fatal issues (e.g. image render unavailable)

    @property
    def has_text(self) -> bool:
        return bool(self.text and self.text.strip())

    @property
    def has_images(self) -> bool:
        return len(self.page_images) > 0


# youtube.com/watch?v=ID, youtu.be/ID, youtube.com/embed/ID, /shorts/ID, /live/ID
_YT_RE = re.compile(
    r"(?:https?://)?(?:www\.)?"
    r"(?:youtube\.com/(?:watch\?(?:[^ \n]*&)?v=|embed/|shorts/|live/)|youtu\.be/)"
    r"([A-Za-z0-9_-]{11})",
    re.IGNORECASE,
)


def find_youtube_url(text: str) -> Optional[str]:
    """Return the first YouTube video URL found in text, normalized to a canonical
    watch URL, or None. The video link is embedded in the slide deck per the
    competition handbook, so this is how Phase 1b gets the video to grade."""
    if not text:
        return None
    m = _YT_RE.search(text)
    if not m:
        return None
    return f"https://www.youtube.com/watch?v={m.group(1)}"


class InputAdapter:
    """Interface: load(source) -> NormalizedInput."""

    def load(self, source: str) -> NormalizedInput:  # pragma: no cover - interface
        raise NotImplementedError


class PDFAdapter(InputAdapter):
    """Load a submission from a PDF (typically a slide deck exported to PDF)."""

    def load(self, source: str) -> NormalizedInput:
        result = NormalizedInput(source=source)
        result.text = self._extract_text(source, result.notes)
        result.video_url = find_youtube_url(result.text)
        result.page_images = self._render_pages(source, result.notes)
        return result

    @staticmethod
    def _extract_text(path: str, notes: List[str]) -> str:
        try:
            import pypdf  # lazy: heavy optional dep
        except ImportError:
            notes.append("pypdf not installed — no text extracted.")
            return ""
        try:
            reader = pypdf.PdfReader(path)
            parts = []
            for page in reader.pages:
                try:
                    parts.append(page.extract_text() or "")
                except Exception as e:  # pragma: no cover - per-page defensive
                    notes.append(f"page text extraction error: {e}")
            return "\n".join(p for p in parts if p)
        except Exception as e:
            notes.append(f"PDF read error: {e}")
            return ""

    @staticmethod
    def _render_pages(path: str, notes: List[str], dpi: int = 130, max_pages: int = 30) -> List[bytes]:
        """Render each PDF page to PNG bytes for the vision model. Slide-deck PDFs
        carry the financials/license/bank/photos as images, so text alone is not
        enough (eng review). Delegates to render_pdf_images (pypdfium2)."""
        images = render_pdf_images(path, dpi=dpi, max_pages=max_pages)
        if not images:
            notes.append("Slide images unavailable (install pypdfium2 + pillow) — "
                         "vision grading will have no slide images.")
        return images


def render_pdf_images(path: str, dpi: int = 130, max_pages: int = 30) -> List[bytes]:
    """Render a PDF's pages to PNG bytes (for vision grading), capped at max_pages
    to bound the multimodal payload. Uses pypdfium2 (Apache-2.0 / BSD-3, Google's
    PDFium) + Pillow — both permissively licensed (no AGPL). Returns [] if the
    libraries are unavailable or the file can't be rendered."""
    try:
        import io
        import pypdfium2 as pdfium
    except ImportError:
        return []
    images: List[bytes] = []
    try:
        pdf = pdfium.PdfDocument(path)
        for i in range(min(len(pdf), max_pages)):
            pil = pdf[i].render(scale=dpi / 72.0).to_pil()
            buf = io.BytesIO()
            pil.save(buf, format="PNG")
            images.append(buf.getvalue())
    except Exception:  # pragma: no cover - defensive (corrupt/encrypted PDF)
        return images
    return images


# --- PowerPoint (.pptx) ----------------------------------------------------- #
def _walk_shapes(shapes):
    """Yield every shape, recursing into group shapes (which hold no text/table of
    their own but contain nested shapes)."""
    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError:  # pragma: no cover - guarded by callers
        MSO_SHAPE_TYPE = None
    for shape in shapes:
        if MSO_SHAPE_TYPE is not None and shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _walk_shapes(shape.shapes)
        else:
            yield shape


def extract_pptx_text(path: str, notes: Optional[List[str]] = None) -> str:
    """Extract text from a .pptx: per-slide shape text + table cells (joined ' | '),
    each slide prefixed '[Slide N]' so the grader can cite locations and financials
    that live in tables are captured. Uses python-pptx (MIT; lxml + Pillow — no AGPL).
    Returns '' if python-pptx is unavailable or the file can't be read."""
    try:
        from pptx import Presentation  # lazy: optional dep
    except ImportError:
        if notes is not None:
            notes.append("python-pptx not installed — no text extracted from .pptx.")
        return ""
    try:
        prs = Presentation(path)
    except Exception as e:
        if notes is not None:
            notes.append(f"PPTX read error: {e}")
        return ""
    slides_out: List[str] = []
    for idx, slide in enumerate(prs.slides, 1):
        parts: List[str] = []
        for shape in _walk_shapes(slide.shapes):
            try:
                if getattr(shape, "has_table", False) and shape.has_table:
                    for row in shape.table.rows:
                        cells = [(c.text or "").strip() for c in row.cells]
                        if any(cells):
                            parts.append(" | ".join(cells))
                elif getattr(shape, "has_text_frame", False) and shape.has_text_frame:
                    t = (shape.text_frame.text or "").strip()
                    if t:
                        parts.append(t)
            except Exception:  # pragma: no cover - defensive per-shape
                continue
        if parts:
            slides_out.append(f"[Slide {idx}]\n" + "\n".join(parts))
    return "\n\n".join(slides_out)


def extract_pptx_images(path: str, notes: Optional[List[str]] = None,
                        max_images: int = 30) -> List[bytes]:
    """Extract embedded picture blobs from a .pptx as PNG bytes for the vision model.
    Full text-on-canvas slide rendering needs LibreOffice; embedded images — photos,
    charts saved as pictures, license/bank scans — are the high-value visual content
    and come free from python-pptx. Blobs are normalized to PNG via Pillow. Returns []
    if python-pptx is unavailable or there are no embedded images."""
    try:
        from pptx import Presentation  # lazy
    except ImportError:
        if notes is not None:
            notes.append("python-pptx not installed — no images extracted from .pptx.")
        return []
    try:
        prs = Presentation(path)
    except Exception:
        return []
    out: List[bytes] = []
    for slide in prs.slides:
        for shape in _walk_shapes(slide.shapes):
            if len(out) >= max_images:
                return out
            try:
                blob = shape.image.blob  # raises AttributeError on non-image shapes
            except Exception:
                continue
            png = _blob_to_png(blob)
            if png:
                out.append(png)
    return out


def _blob_to_png(blob: bytes, max_edge: int = 1600) -> Optional[bytes]:
    """Normalize an image blob to PNG bytes (the vision path serves image/png data
    URIs). Downscales so the long edge is <= max_edge — embedded PPTX photos can be
    huge, and full-res blobs blow up both the encode time and the multimodal payload
    (a PPTX with no LibreOffice render leans on these blobs), making PPTX far slower
    than a fixed-DPI PDF. Falls back to passing through blobs that are already PNG if
    Pillow is absent; drops anything it can't handle."""
    try:
        import io
        from PIL import Image
    except ImportError:
        return blob if blob[:8] == b"\x89PNG\r\n\x1a\n" else None
    try:
        im = Image.open(io.BytesIO(blob)).convert("RGB")
        w, h = im.size
        if max(w, h) > max_edge:
            scale = max_edge / float(max(w, h))
            im = im.resize((max(1, int(w * scale)), max(1, int(h * scale))))
        buf = io.BytesIO()
        im.save(buf, format="PNG")
        return buf.getvalue()
    except Exception:  # pragma: no cover - defensive (unsupported/corrupt image)
        return None


class PPTXAdapter(InputAdapter):
    """Load a submission from a PowerPoint .pptx (text + tables, plus embedded
    images for the vision path)."""

    def load(self, source: str) -> NormalizedInput:
        result = NormalizedInput(source=source)
        result.text = extract_pptx_text(source, result.notes)
        result.video_url = find_youtube_url(result.text)
        result.page_images = extract_pptx_images(source, result.notes)
        if not result.page_images:
            result.notes.append("No embedded images in .pptx — vision grading has no "
                                "slides to see; grade via the text path instead.")
        return result


# --- Word (.docx) ----------------------------------------------------------- #
def extract_docx_text(path: str, notes: Optional[List[str]] = None) -> str:
    """Extract text from a .docx via docx2txt (already the app's DOCX reader, so
    the vision and text paths see the same content).

    docx2txt yields body paragraphs AND table cells, one cell per line. Verified
    against the real submission corpus: every figure in a financial table comes
    through ("$2,400", "$17,000", "Rent", "Salaries"). Row structure is lost — the
    numbers and their labels are not, which is what the grader reads. Returns ''
    if docx2txt is unavailable or the file can't be read."""
    try:
        import docx2txt  # lazy: optional dep
    except ImportError:
        if notes is not None:
            notes.append("docx2txt not installed — no text extracted from .docx.")
        return ""
    try:
        return (docx2txt.process(path) or "").strip()
    except Exception as e:
        if notes is not None:
            notes.append(f"DOCX read error: {e}")
        return ""


def extract_docx_images(path: str, notes: Optional[List[str]] = None,
                        max_images: int = 30) -> List[bytes]:
    """Extract embedded pictures from a .docx as PNG bytes for the vision model.

    Same rationale as the .pptx path: a scanned business licence, a bank statement
    photo, or a product shot is evidence the text cannot carry. A .docx is a zip,
    so its pictures live in word/media/ — no dependency beyond the stdlib. (Most
    prose plans carry none, which is fine: the vision path no longer requires
    images, it just uses them when they exist.)"""
    import zipfile
    out: List[bytes] = []
    try:
        with zipfile.ZipFile(path) as z:
            for name in sorted(n for n in z.namelist() if n.startswith("word/media/")):
                if len(out) >= max_images:
                    break
                try:
                    blob = z.read(name)
                except Exception:  # pragma: no cover - defensive per-entry
                    continue
                if not blob:
                    continue          # directory entry, or an empty part
                png = _blob_to_png(blob)
                if png:
                    out.append(png)
    except Exception as e:
        if notes is not None:
            notes.append(f"DOCX image extraction error: {e}")
    return out


class DOCXAdapter(InputAdapter):
    """Load a submission from a Word .docx (text + table cells, plus any embedded
    images for the vision path). Page rendering would need LibreOffice; the text
    and the embedded pictures are the gradeable content without it."""

    def load(self, source: str) -> NormalizedInput:
        result = NormalizedInput(source=source)
        result.text = extract_docx_text(source, result.notes)
        result.video_url = find_youtube_url(result.text)
        result.page_images = extract_docx_images(source, result.notes)
        return result


# --- Plain text / Markdown (.txt, .md) -------------------------------------- #
def extract_plain_text(path: str, notes: Optional[List[str]] = None) -> str:
    """Read a .txt / .md plan.

    Decoded permissively: an entrant's upload can be cp1252 or carry a stray byte,
    and losing one character beats losing the plan."""
    try:
        with open(path, "rb") as fh:
            raw = fh.read()
    except OSError as e:
        if notes is not None:
            notes.append(f"Could not read text file: {e}")
        return ""
    for enc in ("utf-8-sig", "utf-8", "cp1252"):
        try:
            return raw.decode(enc).strip()
        except UnicodeDecodeError:
            continue
    return raw.decode("utf-8", "replace").strip()


class TextAdapter(InputAdapter):
    """Load a plain-text or Markdown plan.

    There is nothing to render, so page_images is always empty and the vision path
    grades it from its text (see vision_grade.grade_with_vision, which requires
    text OR images rather than images)."""

    def load(self, source: str) -> NormalizedInput:
        result = NormalizedInput(source=source)
        result.text = extract_plain_text(source, result.notes)
        result.video_url = find_youtube_url(result.text)
        if not result.text:
            result.notes.append("No text in the file — nothing to grade.")
        return result


# Every document type get_adapter can load. Single source of truth: the API's
# upload validation and the frontend's file pickers are kept in step with this.
ADAPTER_SUFFIXES = (".pdf", ".pptx", ".docx", ".txt", ".md", ".markdown")


def get_adapter(source: str) -> InputAdapter:
    """Select an adapter for a submission source. PDF, PPTX, DOCX and plain
    text/Markdown are supported; the Google Slides path is a deliberate,
    clearly-signposted extension point."""
    s = source.lower()
    if s.endswith(".pdf"):
        return PDFAdapter()
    if s.endswith(".pptx"):
        return PPTXAdapter()
    if s.endswith(".docx"):
        return DOCXAdapter()
    if s.endswith((".txt", ".md", ".markdown")):
        return TextAdapter()
    if "docs.google.com/presentation" in s or "slides.google.com" in s:
        raise NotImplementedError(
            "Google Slides adapter is future work; export to PDF or PPTX for now."
        )
    raise ValueError(f"No input adapter for source: {source}")
