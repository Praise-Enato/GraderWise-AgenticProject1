"""Tests for PowerPoint (.pptx) submission support.

Authors real .pptx decks in-memory with python-pptx, then round-trips them through
the extractors + adapter. Skipped cleanly if python-pptx is not installed.
"""
import io

import pytest

pptx = pytest.importorskip("pptx")
from pptx import Presentation  # noqa: E402
from pptx.util import Inches  # noqa: E402

from backend.src import input_adapter as IA  # noqa: E402


def _deck_with_text_and_table(path):
    prs = Presentation()
    blank = prs.slide_layouts[6]
    # Slide 1: a title-ish textbox
    s1 = prs.slides.add_slide(blank)
    tb = s1.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
    tb.text_frame.text = "FrostFlow Ice Ventures"
    # Slide 2: a financials table
    s2 = prs.slides.add_slide(blank)
    rows, cols = 2, 2
    tbl = s2.shapes.add_table(rows, cols, Inches(1), Inches(1), Inches(4), Inches(1)).table
    tbl.cell(0, 0).text = "Revenue"
    tbl.cell(0, 1).text = "50000"
    tbl.cell(1, 0).text = "Profit"
    tbl.cell(1, 1).text = "12000"
    prs.save(path)


def _png_bytes():
    from PIL import Image
    buf = io.BytesIO()
    Image.new("RGB", (16, 16), (10, 20, 30)).save(buf, format="PNG")
    return buf.getvalue()


def test_extract_pptx_text_has_slide_markers_text_and_table(tmp_path):
    p = str(tmp_path / "plan.pptx")
    _deck_with_text_and_table(p)
    text = IA.extract_pptx_text(p)
    assert "[Slide 1]" in text
    assert "[Slide 2]" in text
    assert "FrostFlow Ice Ventures" in text
    # table cells joined by ' | ' so financials survive extraction
    assert "Revenue | 50000" in text
    assert "Profit | 12000" in text


def test_get_adapter_routes_pptx(tmp_path):
    p = str(tmp_path / "plan.pptx")
    _deck_with_text_and_table(p)
    adapter = IA.get_adapter(p)
    assert isinstance(adapter, IA.PPTXAdapter)
    norm = adapter.load(p)
    assert norm.has_text
    assert "Revenue | 50000" in norm.text


def test_pptx_adapter_notes_when_no_images(tmp_path):
    p = str(tmp_path / "plan.pptx")
    _deck_with_text_and_table(p)  # no embedded pictures
    norm = IA.get_adapter(p).load(p)
    assert not norm.has_images
    assert any("no slides to see" in n.lower() or "no embedded images" in n.lower()
               for n in norm.notes)


def test_extract_pptx_images_returns_png(tmp_path):
    img_path = str(tmp_path / "pic.png")
    with open(img_path, "wb") as fh:
        fh.write(_png_bytes())
    p = str(tmp_path / "deck.pptx")
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(img_path, Inches(1), Inches(1), Inches(2), Inches(2))
    prs.save(p)

    images = IA.extract_pptx_images(p)
    assert len(images) == 1
    assert images[0][:8] == b"\x89PNG\r\n\x1a\n"  # normalized to real PNG


def test_extract_pptx_text_bad_file_is_graceful(tmp_path):
    p = str(tmp_path / "not-a-deck.pptx")
    with open(p, "wb") as fh:
        fh.write(b"this is not a zip/pptx")
    notes = []
    assert IA.extract_pptx_text(p, notes) == ""
    assert notes  # recorded a read error, did not raise
